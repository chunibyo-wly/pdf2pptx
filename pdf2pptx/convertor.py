import sys
import PIL
import io, os


def _get_pptx_path(pptx_name):
    """
    Return the path to the built-in default .pptx package.
    """
    return os.path.join(os.path.split(__file__)[0], pptx_name)


def pdf2jpeg(pdf_path):
    from pdf2image import convert_from_path, convert_from_bytes

    images = convert_from_path(
        pdf_path,
        dpi=1200,
        fmt="jpeg",
        thread_count=4
    )
    return images


def jpeg2pptx(images, output_path):
    from pptx import Presentation
    from pptx.util import Inches

    ratio = images[0].width / images[0].height
    presentation = None
    if 1.3 <= ratio <= 1.4:
        presentation = Presentation(_get_pptx_path("43.pptx"))
    elif 1.7 <= ratio <= 1.8:
        presentation = Presentation(_get_pptx_path("169.pptx"))
    else:
        presentation = Presentation()

    left = top = Inches(0)
    for image in images:
        # PIL io stream to pptx picture
        with io.BytesIO() as output:
            image.save(output, format="jpeg")
            # add blank layout
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            # image fit size
            slide.shapes.add_picture(output, left, top, presentation.slide_width, presentation.slide_height)

    presentation.save(output_path)
