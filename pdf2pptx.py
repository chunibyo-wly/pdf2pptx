import sys
import PIL
import io


def pdf2jpeg(pdf_path):
    from pdf2image import convert_from_path, convert_from_bytes

    images = convert_from_path(
        pdf_path,
        dpi=1200,
        fmt="jpeg",
        thread_count=4
    )
    return images


def jpeg2pptx(images, output_path):
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()

    left = top = Inches(0)
    for image in images:
        # PIL io stream to pptx picture
        with io.BytesIO() as output:
            image.save(output, format="jpeg")
            # add blank layout
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            # image fit size
            slide.shapes.add_picture(output, left, top, presentation.slide_width, presentation.slide_height)

    presentation.save(output_path)


def main():
    import sys
    images = pdf2jpeg(sys.argv[1])
    jpeg2pptx(images, sys.argv[2])


if __name__ == '__main__':
    main()
